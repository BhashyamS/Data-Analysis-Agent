"""PDF and PowerPoint exports for Analytics Copilot reports."""

from __future__ import annotations

from io import BytesIO
from typing import Any

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


def _plain_text(value: Any) -> str:
    return str(value or "").replace("**", "").replace("##", "").replace("#", "").strip()


def _split_narrative(narrative: str) -> list[str]:
    return [
        line.strip("- ").strip()
        for line in narrative.splitlines()
        if line.strip() and not line.strip().startswith("#")
    ]


def _format_kpi(widget: dict[str, Any]) -> str:
    config = widget.get("config", {})
    value = config.get("value")

    if value is None:
        display = "Not available"
    elif isinstance(value, (int, float)):
        display = f"{value:,.2f}"
    else:
        display = str(value)

    return f"{config.get('prefix', '')}{display}{config.get('suffix', '')}"


def _prepare_chart_data(
    df: pd.DataFrame,
    config: dict[str, Any],
) -> tuple[pd.DataFrame, str, str | None]:
    chart_type = str(config.get("chart_type", "Bar"))
    x_column = str(config.get("x_column", ""))
    y_column = config.get("y_column")
    aggregation = str(config.get("aggregation", "None"))
    top_n = int(config.get("top_n", 20))

    if x_column not in df.columns:
        raise ValueError(f"Column not found: {x_column}")

    working = df.copy()

    if y_column is not None and y_column in working.columns:
        working[y_column] = pd.to_numeric(working[y_column], errors="coerce")

    if aggregation != "None" and y_column:
        aggregation_name = aggregation.lower()

        if aggregation_name == "count":
            working = (
                working.groupby(x_column, dropna=False)[y_column]
                .count()
                .reset_index(name=y_column)
            )
        else:
            working = (
                working.groupby(x_column, dropna=False)[y_column]
                .agg(aggregation_name)
                .reset_index()
            )

        working = working.sort_values(y_column, ascending=False).head(top_n)

    elif chart_type in {"Bar", "Pie"} and y_column is None:
        counts = (
            working[x_column]
            .fillna("Missing")
            .astype(str)
            .value_counts()
            .head(top_n)
        )
        working = counts.rename_axis(x_column).reset_index(name="count")
        y_column = "count"

    return working, x_column, y_column


def chart_widget_to_png(
    df: pd.DataFrame,
    widget: dict[str, Any],
) -> bytes:
    """Render one dashboard chart widget as PNG bytes."""
    config = widget.get("config", {})
    chart_type = str(config.get("chart_type", "Bar"))
    title = str(widget.get("title", "Chart"))
    color_column = config.get("color_column")

    working, x_column, y_column = _prepare_chart_data(df, config)

    figure, axis = plt.subplots(figsize=(9, 5.2))

    if chart_type == "Bar":
        axis.bar(
            working[x_column].astype(str),
            working[y_column] if y_column else range(len(working)),
        )
        axis.tick_params(axis="x", labelrotation=35)

    elif chart_type == "Line":
        dates = pd.to_datetime(working[x_column], errors="coerce")
        axis.plot(dates, working[y_column], marker="o")

    elif chart_type == "Scatter":
        x_values = pd.to_numeric(working[x_column], errors="coerce")
        y_values = pd.to_numeric(working[y_column], errors="coerce")

        if color_column and color_column in working.columns:
            categories = working[color_column].fillna("Missing").astype(str)
            for category in categories.unique():
                mask = categories == category
                axis.scatter(x_values[mask], y_values[mask], label=category, alpha=0.7)
            axis.legend()
        else:
            axis.scatter(x_values, y_values, alpha=0.7)

    elif chart_type == "Pie":
        axis.pie(
            working[y_column],
            labels=working[x_column].astype(str),
            autopct="%1.1f%%",
        )

    elif chart_type == "Histogram":
        values = pd.to_numeric(working[x_column], errors="coerce").dropna()
        axis.hist(values, bins=20)

    elif chart_type == "Box":
        column = y_column or x_column
        values = pd.to_numeric(working[column], errors="coerce").dropna()
        axis.boxplot(values, vert=True)
        axis.set_xticks([1], [column])

    else:
        raise ValueError(f"Unsupported chart type: {chart_type}")

    axis.set_title(title)
    axis.set_xlabel(x_column)

    if y_column and chart_type not in {"Pie", "Box"}:
        axis.set_ylabel(str(y_column))

    figure.tight_layout()

    output = BytesIO()
    figure.savefig(output, format="png", dpi=160, bbox_inches="tight")
    plt.close(figure)
    output.seek(0)
    return output.getvalue()


def build_pdf_export(
    report_data: dict[str, Any],
    df: pd.DataFrame,
    *,
    narrative: str = "",
) -> bytes:
    """Create a complete PDF report."""
    output = BytesIO()
    document = SimpleDocTemplate(
        output,
        pagesize=letter,
        rightMargin=42,
        leftMargin=42,
        topMargin=42,
        bottomMargin=42,
        title=report_data["metadata"]["title"],
    )

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="ReportTitle",
            parent=styles["Title"],
            alignment=TA_CENTER,
            fontSize=24,
            leading=29,
            spaceAfter=18,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Section",
            parent=styles["Heading2"],
            fontSize=16,
            leading=20,
            spaceBefore=15,
            spaceAfter=8,
        )
    )

    story = []
    metadata = report_data["metadata"]
    dataset = report_data["dataset"]
    preparation = report_data["preparation"]
    dashboard = report_data["dashboard"]
    analysis = report_data["analysis"]
    selected = set(metadata["selected_sections"])

    story.append(Paragraph(metadata["title"], styles["ReportTitle"]))
    story.append(
        Paragraph(
            f"Audience: {metadata['audience']}<br/>"
            f"Detail level: {metadata['detail_level']}<br/>"
            f"Generated: {metadata['generated_at']}<br/>"
            f"Source file: {metadata.get('filename') or 'Not provided'}",
            styles["Normal"],
        )
    )
    story.append(Spacer(1, 0.25 * inch))

    if "Executive Summary" in selected:
        story.append(Paragraph("Executive Summary", styles["Section"]))
        summary = narrative.strip() or (
            f"The dataset contains {dataset['rows']:,} rows and "
            f"{dataset['columns']:,} columns, with "
            f"{dataset['missing_percentage']:.2f}% missing cells and "
            f"{dataset['duplicate_rows']:,} duplicate rows."
        )
        story.append(Paragraph(_plain_text(summary), styles["BodyText"]))

    if "Dataset Overview" in selected:
        story.append(Paragraph("Dataset Overview", styles["Section"]))
        table_data = [
            ["Metric", "Value"],
            ["Rows", f"{dataset['rows']:,}"],
            ["Columns", f"{dataset['columns']:,}"],
            ["Missing cells", f"{dataset['missing_cells']:,}"],
            ["Missing percentage", f"{dataset['missing_percentage']:.2f}%"],
            ["Duplicate rows", f"{dataset['duplicate_rows']:,}"],
            ["Memory usage", f"{dataset['memory_mb']:.3f} MB"],
        ]
        table = Table(table_data, colWidths=[2.3 * inch, 3.4 * inch])
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8ECF5")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B8C0D2")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("PADDING", (0, 0), (-1, -1), 7),
                ]
            )
        )
        story.append(table)

    if "Data Preparation" in selected:
        story.append(Paragraph("Data Preparation", styles["Section"]))
        history = preparation.get("history", [])
        if history:
            for item in history:
                description = item.get("description") or item.get("action") or str(item)
                story.append(Paragraph(f"- {_plain_text(description)}", styles["BodyText"]))
        else:
            story.append(Paragraph("No preparation history was recorded.", styles["BodyText"]))

    if "Exploratory Analysis" in selected:
        story.append(Paragraph("Exploratory Analysis", styles["Section"]))
        numeric = analysis.get("numeric_summary", [])[:10]

        if numeric:
            rows = [["Column", "Count", "Missing", "Mean", "Median"]]
            for item in numeric:
                rows.append(
                    [
                        str(item.get("column", "")),
                        str(item.get("count", "")),
                        str(item.get("missing", "")),
                        str(item.get("mean", "")),
                        str(item.get("median", "")),
                    ]
                )

            table = Table(rows, repeatRows=1)
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8ECF5")),
                        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#B8C0D2")),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("PADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            story.append(table)
        else:
            story.append(Paragraph("No numeric summary was available.", styles["BodyText"]))

    if "Statistical Analysis" in selected:
        story.append(Paragraph("Statistical Analysis", styles["Section"]))
        relationships = analysis.get("strongest_correlations", [])[:10]

        if relationships:
            rows = [["Variable 1", "Variable 2", "Correlation", "Strength"]]
            for item in relationships:
                rows.append(
                    [
                        str(item.get("variable_1", "")),
                        str(item.get("variable_2", "")),
                        str(item.get("correlation", "")),
                        str(item.get("strength", "")),
                    ]
                )

            table = Table(rows, repeatRows=1)
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8ECF5")),
                        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#B8C0D2")),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("PADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            story.append(table)
        else:
            story.append(Paragraph("No qualifying relationships were detected.", styles["BodyText"]))

    if "AI Insights" in selected:
        story.append(Paragraph("AI Insights", styles["Section"]))
        ai_report = report_data.get("ai_report", "")
        story.append(
            Paragraph(
                _plain_text(ai_report) if ai_report else "No saved AI report was available.",
                styles["BodyText"],
            )
        )

    if "Dashboard" in selected:
        story.append(PageBreak())
        story.append(Paragraph(dashboard["title"], styles["ReportTitle"]))

        for widget in dashboard.get("widgets", []):
            widget_type = widget.get("type")
            title = str(widget.get("title", "Widget"))
            config = widget.get("config", {})

            if widget_type == "kpi":
                story.append(
                    KeepTogether(
                        [
                            Paragraph(title, styles["Heading3"]),
                            Paragraph(_format_kpi(widget), styles["Heading2"]),
                            Spacer(1, 0.1 * inch),
                        ]
                    )
                )

            elif widget_type == "chart":
                try:
                    png = chart_widget_to_png(df, widget)
                    story.append(Paragraph(title, styles["Heading3"]))
                    story.append(Image(BytesIO(png), width=6.6 * inch, height=3.8 * inch))
                    story.append(Spacer(1, 0.15 * inch))
                except Exception as exc:
                    story.append(
                        Paragraph(
                            f"{title}: chart could not be rendered ({_plain_text(exc)}).",
                            styles["BodyText"],
                        )
                    )

            else:
                story.append(Paragraph(title, styles["Heading3"]))
                story.append(Paragraph(_plain_text(config.get("text", "")), styles["BodyText"]))

    if "Recommendations" in selected:
        story.append(Paragraph("Recommendations", styles["Section"]))
        narrative_lines = _split_narrative(narrative)
        recommendations = [
            line for line in narrative_lines
            if "recommend" in line.lower() or "action" in line.lower()
        ]

        if recommendations:
            for item in recommendations[:8]:
                story.append(Paragraph(f"- {_plain_text(item)}", styles["BodyText"]))
        else:
            story.append(
                Paragraph(
                    "Review the strongest relationships, outliers, and dashboard KPIs "
                    "before making decisions.",
                    styles["BodyText"],
                )
            )

    if "Limitations" in selected:
        story.append(Paragraph("Limitations", styles["Section"]))
        limitations = [
            "Results describe the uploaded dataset and may not generalize beyond it.",
            "Correlation and group differences do not establish causation.",
            "Missing values, small groups, and outliers may affect findings.",
            "AI-generated text should be reviewed before use in decisions.",
        ]
        for item in limitations:
            story.append(Paragraph(f"- {item}", styles["BodyText"]))

    document.build(story)
    output.seek(0)
    return output.getvalue()


def _add_title(slide, title: str, subtitle: str = "") -> None:
    title_box = slide.shapes.add_textbox(
        Inches(0.65),
        Inches(0.35),
        Inches(12.0),
        Inches(0.65),
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(1.05),
            Inches(11.8),
            Inches(0.45),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(13)


def _add_bullet_slide(
    presentation: Presentation,
    title: str,
    bullets: list[str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, title)

    box = slide.shapes.add_textbox(
        Inches(0.85),
        Inches(1.35),
        Inches(11.6),
        Inches(5.6),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = _plain_text(bullet)
        paragraph.level = 0
        paragraph.font.size = Pt(19)
        paragraph.space_after = Pt(11)


def build_powerpoint_export(
    report_data: dict[str, Any],
    df: pd.DataFrame,
    *,
    narrative: str = "",
) -> bytes:
    """Create a stakeholder-ready PowerPoint report."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    metadata = report_data["metadata"]
    dataset = report_data["dataset"]
    preparation = report_data["preparation"]
    analysis = report_data["analysis"]
    dashboard = report_data["dashboard"]
    selected = set(metadata["selected_sections"])

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(1.4),
    )
    title_frame = title_box.text_frame
    title_frame.text = metadata["title"]
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = title_slide.shapes.add_textbox(
        Inches(1.3),
        Inches(3.6),
        Inches(10.7),
        Inches(1.0),
    )
    subtitle = subtitle_box.text_frame
    subtitle.text = (
        f"{metadata['audience']} report | {metadata['generated_at']}\n"
        f"Source: {metadata.get('filename') or 'Not provided'}"
    )
    subtitle.paragraphs[0].font.size = Pt(17)
    subtitle.paragraphs[0].alignment = PP_ALIGN.CENTER

    if "Executive Summary" in selected:
        summary = narrative.strip() or (
            f"The dataset contains {dataset['rows']:,} rows across "
            f"{dataset['columns']:,} columns. Missing cells represent "
            f"{dataset['missing_percentage']:.2f}% of the dataset."
        )
        _add_bullet_slide(
            presentation,
            "Executive Summary",
            _split_narrative(summary)[:7] or [summary],
        )

    if "Dataset Overview" in selected:
        _add_bullet_slide(
            presentation,
            "Dataset Overview",
            [
                f"Rows: {dataset['rows']:,}",
                f"Columns: {dataset['columns']:,}",
                f"Missing cells: {dataset['missing_cells']:,} "
                f"({dataset['missing_percentage']:.2f}%)",
                f"Duplicate rows: {dataset['duplicate_rows']:,}",
                f"Memory usage: {dataset['memory_mb']:.3f} MB",
            ],
        )

    if "Data Preparation" in selected:
        history = preparation.get("history", [])
        bullets = [
            item.get("description") or item.get("action") or str(item)
            for item in history[:10]
        ]
        _add_bullet_slide(
            presentation,
            "Data Preparation",
            bullets or ["No preparation history was recorded."],
        )

    if "Statistical Analysis" in selected:
        relationships = analysis.get("strongest_correlations", [])[:7]
        bullets = [
            f"{item.get('variable_1')} and {item.get('variable_2')}: "
            f"{item.get('correlation')} ({item.get('strength')})"
            for item in relationships
        ]
        _add_bullet_slide(
            presentation,
            "Strongest Relationships",
            bullets or ["No qualifying correlations were detected."],
        )

    if "AI Insights" in selected and report_data.get("ai_report"):
        _add_bullet_slide(
            presentation,
            "AI Insights",
            _split_narrative(report_data["ai_report"])[:8],
        )

    if "Dashboard" in selected:
        kpis = [
            widget for widget in dashboard.get("widgets", [])
            if widget.get("type") == "kpi"
        ]

        if kpis:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _add_title(slide, dashboard["title"], "Key performance indicators")

            card_width = 3.7
            for index, widget in enumerate(kpis[:3]):
                left = 0.7 + index * 4.15
                shape = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(2.0),
                    Inches(card_width),
                    Inches(2.4),
                )
                frame = shape.text_frame
                frame.clear()

                title_p = frame.paragraphs[0]
                title_p.text = str(widget.get("title", "KPI"))
                title_p.font.size = Pt(17)
                title_p.font.bold = True
                title_p.alignment = PP_ALIGN.CENTER

                value_p = frame.add_paragraph()
                value_p.text = _format_kpi(widget)
                value_p.font.size = Pt(30)
                value_p.font.bold = True
                value_p.alignment = PP_ALIGN.CENTER

        for widget in dashboard.get("widgets", []):
            if widget.get("type") != "chart":
                continue

            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _add_title(slide, str(widget.get("title", "Dashboard Chart")))

            try:
                png = chart_widget_to_png(df, widget)
                image_stream = BytesIO(png)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(0.8),
                    Inches(1.25),
                    width=Inches(11.7),
                    height=Inches(5.7),
                )
            except Exception as exc:
                _add_bullet_slide(
                    presentation,
                    str(widget.get("title", "Dashboard Chart")),
                    [f"Chart could not be rendered: {_plain_text(exc)}"],
                )

        insights = [
            widget for widget in dashboard.get("widgets", [])
            if widget.get("type") in {"insight", "text"}
        ]

        if insights:
            _add_bullet_slide(
                presentation,
                "Dashboard Insights",
                [
                    f"{widget.get('title', 'Insight')}: "
                    f"{widget.get('config', {}).get('text', '')}"
                    for widget in insights[:7]
                ],
            )

    if "Recommendations" in selected:
        lines = _split_narrative(narrative)
        recommendations = [
            line for line in lines
            if "recommend" in line.lower() or "action" in line.lower()
        ]
        _add_bullet_slide(
            presentation,
            "Recommendations",
            recommendations[:7]
            or [
                "Validate the highest-impact findings with domain experts.",
                "Review outliers and small sample groups before taking action.",
                "Track the selected KPIs over time.",
            ],
        )

    if "Limitations" in selected:
        _add_bullet_slide(
            presentation,
            "Limitations",
            [
                "The findings describe the uploaded dataset only.",
                "Correlation does not establish causation.",
                "Missing data, outliers, and small groups may affect results.",
                "AI-generated content should be reviewed before decisions.",
            ],
        )

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output.getvalue()
